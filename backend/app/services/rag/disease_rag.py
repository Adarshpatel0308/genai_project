"""
Disease & Crop Knowledge RAG Pipeline
- Separate ChromaDB collection from KrishiGPT
- Supports PDF, PPT, PPTX, TXT, images
- Used to enhance disease detection with admin-uploaded knowledge
"""
import os
import re
import asyncio
from pathlib import Path
from app.core.config import settings

_embeddings = None
_disease_vectorstore = None
_soil_vectorstore = None


def get_embeddings():
    global _embeddings
    if _embeddings is None:
        try:
            from langchain_huggingface import HuggingFaceEmbeddings
            _embeddings = HuggingFaceEmbeddings(
                model_name="sentence-transformers/paraphrase-multilingual-MiniLM-L12-v2",
                model_kwargs={"device": "cpu"}
            )
        except ImportError:
            from langchain_community.embeddings import FakeEmbeddings
            _embeddings = FakeEmbeddings(size=384)
    return _embeddings


def get_disease_vectorstore():
    global _disease_vectorstore
    if _disease_vectorstore is None:
        from langchain_chroma import Chroma
        _disease_vectorstore = Chroma(
            persist_directory=settings.CHROMA_PERSIST_DIR,
            embedding_function=get_embeddings(),
            collection_name="disease_knowledge"
        )
    return _disease_vectorstore


def get_soil_vectorstore():
    global _soil_vectorstore
    if _soil_vectorstore is None:
        from langchain_chroma import Chroma
        _soil_vectorstore = Chroma(
            persist_directory=settings.CHROMA_PERSIST_DIR,
            embedding_function=get_embeddings(),
            collection_name="soil_crop_knowledge"
        )
    return _soil_vectorstore


def extract_text_from_pdf(file_path: str) -> list[str]:
    """Extract text chunks from PDF."""
    from pypdf import PdfReader
    reader = PdfReader(file_path)
    chunks = []
    for page in reader.pages:
        text = page.extract_text() or ""
        if text.strip():
            chunks.append(text.strip())
    return chunks


def extract_text_from_pptx(file_path: str) -> list[str]:
    """Extract text from each slide of a PPTX file."""
    from pptx import Presentation
    prs = Presentation(file_path)
    chunks = []
    for i, slide in enumerate(prs.slides):
        slide_text = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_text.append(shape.text.strip())
        if slide_text:
            chunks.append(f"Slide {i+1}:\n" + "\n".join(slide_text))
    return chunks


def extract_text_from_txt(file_path: str) -> list[str]:
    with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
        text = f.read()
    # Split into ~800 char chunks
    return [text[i:i+800] for i in range(0, len(text), 700) if text[i:i+800].strip()]


def extract_text_from_image(file_path: str) -> list[str]:
    """Use Tesseract OCR to extract text from image."""
    try:
        import pytesseract
        from PIL import Image
        img = Image.open(file_path)
        text = pytesseract.image_to_string(img, lang="eng+hin")
        if text.strip():
            return [text.strip()]
    except Exception as e:
        print(f"⚠️  Image OCR failed: {e}")
    return []


def extract_chunks(file_path: str) -> list[str]:
    """Extract text chunks from any supported file type."""
    ext = Path(file_path).suffix.lower()
    if ext == ".pdf":
        return extract_text_from_pdf(file_path)
    elif ext in [".ppt", ".pptx"]:
        return extract_text_from_pptx(file_path)
    elif ext == ".txt":
        return extract_text_from_txt(file_path)
    elif ext in [".jpg", ".jpeg", ".png", ".webp", ".bmp"]:
        return extract_text_from_image(file_path)
    return []


def index_disease_document(file_path: str, metadata: dict) -> int:
    """Index a document into the disease knowledge collection."""
    chunks = extract_chunks(file_path)
    if not chunks:
        return 0

    from langchain_core.documents import Document
    docs = [
        Document(page_content=chunk, metadata={**metadata, "source": file_path, "chunk": i})
        for i, chunk in enumerate(chunks)
        if chunk.strip()
    ]

    vs = get_disease_vectorstore()
    vs.add_documents(docs)
    print(f"✅ Indexed {len(docs)} chunks from {Path(file_path).name} → disease_knowledge")
    return len(docs)


def index_soil_document(file_path: str, metadata: dict) -> int:
    """Index a document into the soil/crop knowledge collection."""
    chunks = extract_chunks(file_path)
    if not chunks:
        return 0

    from langchain_core.documents import Document
    docs = [
        Document(page_content=chunk, metadata={**metadata, "source": file_path, "chunk": i})
        for i, chunk in enumerate(chunks)
        if chunk.strip()
    ]

    vs = get_soil_vectorstore()
    vs.add_documents(docs)
    print(f"✅ Indexed {len(docs)} chunks from {Path(file_path).name} → soil_crop_knowledge")
    return len(docs)


def search_disease_knowledge(query: str, k: int = 4) -> str:
    """Search disease knowledge base for relevant information."""
    try:
        vs = get_disease_vectorstore()
        docs = vs.similarity_search(query, k=k)
        if not docs:
            return ""
        return "\n\n---\n\n".join([d.page_content for d in docs])
    except Exception as e:
        print(f"⚠️  Disease knowledge search failed: {e}")
        return ""


def search_soil_knowledge(query: str, k: int = 4) -> str:
    """Search soil/crop knowledge base for relevant information."""
    try:
        vs = get_soil_vectorstore()
        docs = vs.similarity_search(query, k=k)
        if not docs:
            return ""
        return "\n\n---\n\n".join([d.page_content for d in docs])
    except Exception as e:
        print(f"⚠️  Soil knowledge search failed: {e}")
        return ""


async def get_enhanced_disease_context(crop: str, disease: str, language: str = "hi") -> str:
    """
    Search disease knowledge base + generate enhanced explanation using LLM.
    Called after image classification to enrich the treatment plan.
    """
    query = f"{crop} {disease} treatment symptoms prevention"
    kb_context = search_disease_knowledge(query, k=4)

    if not kb_context:
        return ""

    from app.services.ai.llm_service import run_chain
    lang_name = {"hi": "Hindi", "en": "English", "mr": "Marathi", "gu": "Gujarati"}.get(language, "Hindi")

    system = f"""You are an expert plant pathologist for Indian farmers.
Using the knowledge base information provided, give a detailed disease management guide in {lang_name}.
Include: disease description, symptoms, chemical treatment with doses, organic treatment, prevention steps.
Be specific with product names and quantities."""

    prompt = f"""Crop: {crop}
Disease: {disease}

Knowledge Base Information:
{kb_context[:3000]}

Provide a comprehensive treatment and prevention guide."""

    try:
        result = await asyncio.wait_for(run_chain(system, prompt), timeout=25.0)
        return result if result and len(result) > 100 else ""
    except Exception:
        return ""
